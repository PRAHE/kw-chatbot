'''
파일 -> 텍스트 추출
'''
import io
from typing import Tuple, Dict, Any

from pypdf import PdfReader
from pptx import Presentation
from PIL import Image

def extract_text_from_txt_bytes(data: bytes) -> str:
    return data.decode("utf-8", errors="ignore")

def extract_text_from_pdf_bytes(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    texts = []
    for i, page in enumerate(reader.pages):
        t = page.extract_text() or ""
        t = t.strip()
        if t:
            texts.append(f"[PDF p.{i+1}]\n{t}")
    return "\n\n".join(texts).strip()

def extract_text_from_pptx_bytes(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    out = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            # 텍스트
            if hasattr(shape, "text") and shape.text:
                slide_texts.append(shape.text)
            # 표
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                rows = []
                for r in table.rows:
                    rows.append(" | ".join((c.text or "").strip() for c in r.cells))
                if rows:
                    slide_texts.append("[TABLE]\n" + "\n".join(rows))

        joined = "\n".join(t.strip() for t in slide_texts if t.strip())
        if joined:
            out.append(f"[PPTX slide {idx}]\n{joined}")
    return "\n\n".join(out).strip()

def extract_text_from_image_bytes(data: bytes) -> Tuple[str, Dict[str, Any]]:
    # OCR. tesseract가 없으면 빈 텍스트 + 안내 메타 반환
    meta: Dict[str, Any] = {"ocr": "unavailable"}
    try:
        import pytesseract # noqa
        img = Image.open(io.BytesIO(data)).convert("RGB")
        text = pytesseract.image_to_string(img, lang="eng+kor")
        text = (text or "").strip()
        meta["ocr"] = "tesseract"
        return text, meta
    except Exception as e:
        meta["ocr_error"] = str(e)
        return "", meta
    
def extract_text(filename: str, data: bytes) -> Tuple[str, Dict[str, Any]]:
    """
    반환: (text, meta)
    meta에는 extractor/상태 등을 넣어둬서 나중에 디버깅/출처표시에 활용
    """
    lower = filename.lower()
    meta: Dict[str, Any] = {"filename": filename}

    if lower.endswith(".txt"):
        meta["type"] = "txt"
        return extract_text_from_txt_bytes(data), meta
    
    if lower.endswith(".pdf"):
        meta["type"] = "pdf"
        return extract_text_from_pdf_bytes(data), meta
    
    if lower.endswith(".pptx"):
        meta["type"] = "pptx"
        return extract_text_from_pptx_bytes(data), meta
    
    if lower.endswith((".png," ".jpg", ".jpeg", ".webp")):
        meta["type"] = "image"
        text, ocr_meta = extract_text_from_image_bytes(data)
        meta.update(ocr_meta)
        return text, meta
    
    # MVP 정책: hwp/docx 등은 우선 미지원
    meta["type"] = "unsupported"
    return "", meta